import collections.abc
from pptx import Presentation

prs = Presentation()
title_slide_layout = prs.slide_layouts[3]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "変更後のパワーポイントファイル"
subtitle.text = "レイアウトを変更するとどうなる？"

count = 0
for slide in prs.slides:
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                count += len(run.text)

print(f'パワーポイント内の文字数は{count}個です。')

prs.save('test2.pptx')